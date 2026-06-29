"""Generate REFINEMENTS_SLIDES.pptx from the refinement-cycle review content.

One-off generator. Run:  python generate_slides.py
Requires: python-pptx  (pip install python-pptx)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Theme ─────────────────────────────────────────────────────────────────
BG = RGBColor(0x0A, 0x0A, 0x0A)        # near-black, matches app
WHITE = RGBColor(0xF5, 0xF5, 0xF5)
MUTED = RGBColor(0xB0, 0xB8, 0xC0)
ACCENT = RGBColor(0x4D, 0xD0, 0xE1)    # cyan accent
GOOD = RGBColor(0x66, 0xBB, 0x6A)      # green for "result"

WIDE_W, WIDE_H = Inches(13.333), Inches(7.5)


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(
        1, Inches(0.0), Inches(0.0), Inches(0.18), WIDE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False


def set_runs(paragraph, text, color, size, bold_default=False):
    """Render text with **bold** inline markup into runs."""
    parts = text.split("**")
    for i, part in enumerate(parts):
        if part == "":
            continue
        run = paragraph.add_run()
        run.text = part
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold_default or (i % 2 == 1)
        run.font.name = "Segoe UI"


def title_slide(prs, title, subtitle, presenter):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Segoe UI"

    p2 = tf.add_paragraph(); p2.space_before = Pt(14)
    r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(26); r2.font.color.rgb = ACCENT; r2.font.name = "Segoe UI"

    p3 = tf.add_paragraph(); p3.space_before = Pt(24)
    r3 = p3.add_run(); r3.text = presenter
    r3.font.size = Pt(16); r3.font.color.rgb = MUTED; r3.font.name = "Segoe UI"
    return slide


def content_slide(prs, title, bullets):
    """bullets: list of (level:int, text:str, kind:str) where kind in {'', 'good', 'muted'}."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Segoe UI"

    # Underline accent
    ln = slide.shapes.add_shape(1, Inches(0.72), Inches(1.45), Inches(3.2), Inches(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background(); ln.shadow.inherit = False

    # Body
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.75), Inches(12.0), Inches(5.3))
    bf = body.text_frame; bf.word_wrap = True
    first = True
    for level, text, kind in bullets:
        p = bf.paragraphs[0] if first else bf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(7)
        color = GOOD if kind == "good" else (MUTED if kind == "muted" else WHITE)
        size = 19 if level == 0 else 16
        bullet = "•  " if level == 0 else "–  "
        if kind == "none":
            bullet = ""
        set_runs(p, bullet + text, color, size)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H

    title_slide(
        prs,
        "Buddy Walk",
        "Refinement & Accessibility Hardening",
        "Focus: refine what exists — especially for blind & low-vision users",
    )

    content_slide(prs, "Why This Cycle", [
        (0, "Core features already exist: Q&A, camera describe, voice input, haptic navigation, companion sharing, saved places.", ""),
        (0, "New features are **on the back burner.**", ""),
        (0, "Goal: make what we have **reliable, clear, and comfortable** to use by ear and touch.", ""),
        (0, 'Every fix judged by one question:', ""),
        (1, '"What does this feel like if you can\'t see the screen?"', "muted"),
    ])

    content_slide(prs, "How We Approached It", [
        (0, "**Audited** the full user path read-only before changing anything.", ""),
        (0, "Ranked findings by **impact on a non-sighted user**, not by code difficulty.", ""),
        (0, "Fixed highest-impact items first; verified types + lint after each.", ""),
    ])

    content_slide(prs, "Audit Findings (Snapshot)", [
        (0, "Critical", "good"),
        (1, "Double narration — app TTS and screen reader both read the same text.", ""),
        (1, "Azure speech token never refreshed — voice dies mid-session.", ""),
        (1, "Pop-up modals after every photo/video — friction for a blind user.", ""),
        (0, "High", "good"),
        (1, "Shake-to-listen misfires while walking.", ""),
        (1, "Chat logs recorded empty questions.", ""),
        (0, "Medium / backlog", "good"),
        (1, "Long silent waits, no off-route reroute, dual GPS watchers.", ""),
    ])

    content_slide(prs, "Fix #1: One Voice, Not Two  (biggest win)", [
        (0, "**Problem:** with a screen reader on, the app's TTS and the OS reader both spoke — overlapping, doubled narration everywhere.", ""),
        (0, "**Fix:** new announce() utility detects the screen reader.", ""),
        (1, "Screen reader ON → let the screen reader speak.", ""),
        (1, "Screen reader OFF → use the app's own TTS.", ""),
        (1, "Removed duplicate speak+announce pairs and competing live-regions.", ""),
        (0, "**Result:** exactly one clear narration, every time.", "good"),
    ])

    content_slide(prs, "Fix #2: Voice Input That Survives Long Walks", [
        (0, "**Problem:** token fetched once at launch; Azure tokens expire (~10 min), so voice input quietly failed on longer walks until restart.", ""),
        (0, "**Fix:**", ""),
        (1, "Refresh the token automatically every ~8 minutes.", ""),
        (1, "On an unauthorized transcription, refresh once and retry transparently.", ""),
        (0, "**Result:** voice input keeps working for the whole outing.", "good"),
    ])

    content_slide(prs, "Fix #3: No More Interrupting Pop-ups", [
        (0, "**Problem:** a native modal appeared after every photo/video; a blind user had to find and dismiss an OK button before continuing.", ""),
        (0, "**Fix:** removed the modals; kept the spoken confirmation.", ""),
        (0, "**Result:** capture → continue immediately. No hunting for a button.", "good"),
    ])

    content_slide(prs, "Fix #4: Shake-to-Listen That Knows You're Walking", [
        (0, "**Problem:** the shake detector tripped too easily — a normal stride could open the mic.", ""),
        (0, "**Fix:**", ""),
        (1, "Raised threshold 1.8 → 2.6 G (a stride peaks below a deliberate shake).", ""),
        (1, "Require two strong spikes within 700 ms, so one footfall won't trip it.", ""),
        (0, "**Result:** the mic opens when you mean it to — not when you take a step.", "good"),
    ])

    content_slide(prs, "Fix #5: Honest Chat Logs", [
        (0, "**Problem:** the question was cleared before logging, so history often saved an empty question next to the answer.", ""),
        (0, "**Fix:** snapshot the question the moment it's sent, then log that.", ""),
        (0, "**Result:** history is actually useful for review and debugging.", "good"),
    ])

    content_slide(prs, "Fix #6: More Concise AI Responses", [
        (0, "**Problem:** answers were long and meandering — tiring when every word is read aloud. The prompt even said \"give multiple options.\"", ""),
        (0, "**Fix (server prompt rewrite):**", ""),
        (1, "Default to 1–2 short sentences, lead with the answer, no filler.", ""),
        (1, "Single best answer; list options (max 3) only when explicitly asked.", ""),
        (1, "Kept all safety rules: provided data only, no coords, turn left/right phrasing.", ""),
        (1, "No hard token cap — would truncate long turn-by-turn directions.", "muted"),
    ])

    content_slide(prs, "Fix #7: Feedback While You Wait", [
        (0, "**Problem:** on a slow response a blind user got only a spinner — silence for up to a minute.", ""),
        (0, "**Fix:**", ""),
        (1, 'Spoken reassurance after ~9 s: "Still working on your request."', ""),
        (1, "Distinguish a timeout (\"server may be busy\") from offline (\"no internet\").", ""),
        (0, "**Result:** no more misleading \"you're offline\" when you're not.", "good"),
    ])

    content_slide(prs, "Fix #8: Stop Talking When You Leave", [
        (0, "**Problem:** the AI answer kept speaking after opening Navigation, Companion, or Saved Places — talking over the next screen.", ""),
        (0, "**Fix:** speech stops the moment you navigate away.", ""),
        (0, "**Result:** haptic navigation cues start cleanly.", "good"),
    ])

    content_slide(prs, "Context: Earlier Work This Built On", [
        (0, "**Stronger haptics:** unified tap/maneuver vibration vocabulary on every button.", ""),
        (0, "**Bloated-navigation guard:** stopped \"Brooklyn → out-of-state → Brooklyn\" routes via tighter geocoding, a distance sanity check, hardened prompts, and a client-side >12 km guard that suggests transit.", ""),
        (0, "**Audio session race fix:** resolved the recurring iOS audio error (OSStatus 561017449).", ""),
        (0, "**Auth bypass:** drops straight into the app for now (sign-up gating deferred).", ""),
    ])

    content_slide(prs, "Note for the Team: Where the Backend Lives", [
        (0, "App's default backend is **buddywalk.app** (hosted upstream — we don't control it).", ""),
        (0, "The **concise-response prompt is server-side** — it only takes effect when the app talks to a server running our code.", ""),
        (0, "To see it live: run our mobile/server and point the app via EXPO_PUBLIC_API_URL, or deploy our backend.", ""),
        (0, "**All mobile-side fixes (#1, #3, #4, #5, #7, #8) are already live — no deploy needed.**", "good"),
    ])

    content_slide(prs, "How to Test (Quick Script)", [
        (0, "Screen reader ON → ask a question → answer is spoken **once.**", ""),
        (0, "Take a photo → **no pop-up**, just spoken confirmation.", ""),
        (0, "Walk normally for a minute → mic should **not** open on its own.", ""),
        (0, "Use voice after 10+ minutes → it should **still work** (token refresh).", ""),
        (0, "Ask on a poor connection → hear \"still working,\" then a clear timeout vs offline message.", ""),
        (0, "Open Navigation/Companion mid-answer → speech **stops.**", ""),
    ])

    content_slide(prs, "Backlog / Next Candidates", [
        (0, "**Off-route recovery:** offer \"recalculate from here,\" not just a warning.", ""),
        (0, "**Battery:** avoid two high-accuracy GPS watchers at once during navigation.", ""),
        (0, "**Camera gesture w/ screen reader:** tap-vs-hold is unreliable under VoiceOver.", ""),
        (0, "**Tighten the entrance-description prompt** while keeping safety-critical detail.", ""),
        (0, "**Host our own backend** so server improvements reach users directly.", ""),
    ])

    content_slide(prs, "Summary", [
        (0, "Shifted from building features to **making the experience trustworthy** for blind/low-vision users.", ""),
        (0, "**8 user-facing refinements** shipped on mobile; **1 server prompt** improvement staged.", ""),
        (0, "Everything verified (types + lint), grounded in \"how does this feel without sight?\"", ""),
        (0, "Questions / discussion.", "muted"),
    ])

    out = "REFINEMENTS_SLIDES.pptx"
    prs.save(out)
    print(f"Wrote {out} with {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    build()
